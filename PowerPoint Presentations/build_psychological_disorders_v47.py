# -*- coding: utf-8 -*-
"""Build V4.7 Psychological Disorders and their Assessment from V4.32 base."""
from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

SOURCE = Path(
    r"C:\Users\mike\Documents\PET\PsychoPathology and Assessments"
    r"\PowerPoint Presentations\V4.32 Psychological Disorders and their Assessment.pptx"
)
OUT_DIR = Path(__file__).resolve().parent
OUT = OUT_DIR / "V4.7 Psychological Disorders and their Assessment.pptx"
WORK = OUT_DIR / "_v47_work.pptx"

TITLE_PT = 32
BODY_PT = 22
BODY_FONT = "Arial Narrow"
TITLE_FONT = "Calibri"
NOTES_PT = 11

VERSION_NOTE = (
    "V4.7 revised for tighter pacing, stronger learning-objective alignment, "
    "lighter morning exercises, and retained exercise instructions."
)

EXercises = {
    "purpose_scope": """Title: Purpose and scope of initial assessment
Aim: Clarify what an initial assessment is for and what it should include.
Original time: 20 minutes
Pairs: 15 min discussion + 5 min share-back
Q1: Purpose of initial assessment?
Q2: What should it include?""",
    "rahul": """Title: Client suitability — case 'Rahul'
Aim: Apply suitability criteria and psychosynthesis lens to a complex referral.
Original time: 40 minutes
Three groups with case handout
Assess suitability; identify risks/referral needs
2-minute rationale: work now, wait, or refer""",
    "defence_partner": """Question: What are your go-to defense mechanisms?
Original exercise: 10 min w partner
10-minute partner brainstorm on own defences.""",
    "transference_risk": """Title: Transference, countertransference, and risk
Aim: Practise reading relational dynamics and risk in simulated assessment.
Original time: 45 minutes
Five scenarios in groups
Pairs: identify transference, countertransference, risk
5-minute debrief per scenario""",
    "formulation": """Title: Integrated case formulation workshop
Aim: Produce a psychosynthesis-informed formulation and assessment plan.
Original time: 60 minutes
Groups: Ellen (BPD) or Kim (eating disorder)
Use formulation map; include risk and CORE/PHQ
5-minute group presentation""",
}

# Final deck order: (title match, optional body substring, occurrence for duplicates)
FINAL_ORDER: list[tuple[str, str | None, int]] = [
    ("Psychological Disorders and their Assessment", "Weekend Workshop", 1),
    ("Day 1", "Saturday Morning", 1),
    ("Learning aims", None, 1),
    ("Workshop schedule", None, 1),
    ("Assessing the individual", None, 1),
    ("Many forms of troubling behaviour", None, 1),
    ("Power", "Threat", 1),
    ("Psychological Mindedness", None, 1),
    ("Disidentification as an assessment lens", None, 1),
    ('Observing self / "I"', None, 1),
    ("Psychosynthesis egg diagram as an assessment map", None, 1),
    ("Psychosynthesis assessment questions", None, 1),
    ("Subpersonalities in assessment", None, 1),
    ("Psychosynthesis capacities for therapy", None, 1),
    ("Love and will in assessment", None, 1),
    ("Case study: Chen", None, 1),
    ("Indicators of client suitability", None, 1),
    ("Higher concern / referral indicators", None, 1),
    ("Intersectionality in assessment", None, 1),
    ("Family systems and hidden dynamics", None, 1),
    ("Initial screening tools", None, 1),
    ("CORE-OM", None, 1),
    ("CORE-10", None, 1),
    ("PHQ-9", None, 1),
    ("GAD-7", None, 1),
    ("Session Exercise", "Purpose and scope", 1),
    ("Day 1", "Saturday Afternoon", 1),
    ("Online vs in-person assessment", None, 1),
    ("Assessing for suicide and self-harm", None, 1),
    ("Defense mechanisms", "QUESTIONS", 1),
    ("Defence mechanisms 1", None, 1),
    ("Defence mechanisms 6", None, 1),
    ("Mental status exam", None, 1),
    ("MSE in practice", None, 1),
    ("Psychosynthesis assessment domains", None, 1),
    ("Session Exercise", "Rahul", 1),
    ("Psychological Disorders and their Assessment", "Sunday", 1),
    ("Day 2", "Sunday Morning", 1),
    ("Opening reflection", None, 1),
    ("Brief reflection (optional)", None, 1),
    ("transference", None, 1),
    ("counter-transference", None, 1),
    ("Transferential field", None, 1),
    ("Dissociation in the transferential field", None, 1),
    ("Risk Assessment", None, 1),
    ("What questions do I ask", None, 1),
    ("Self-harm vs suicidality", None, 1),
    ("Dissociation — screening", None, 1),
    ("Session Exercise", "Transference, countertransference", 1),
    ("Day 2", "Sunday Afternoon", 1),
    ("Understanding mild to moderate", None, 1),
    ("Neurosis, psychosis", None, 1),
    ("Lower, middle, and higher unconscious", None, 1),
    ("Personality disorders", None, 1),
    ("Borderline patterns", None, 1),
    ("Narcissistic patterns", None, 1),
    ("Writing session notes", None, 1),
    ("ICD-11 / DSM-5", None, 1),
    ("Eating disorders", None, 1),
    ("From hypothesis to review", None, 1),
    ("Hypothesising and formulation", None, 1),
    ("Psychosynthesis formulation map", None, 1),
    ("What makes an assessment psychosynthesis-informed", None, 1),
    ("Cognitive tools and negative automatic thoughts", None, 1),
    ("Psychosynthesis and other models", None, 1),
    ("Session Exercise", "Integrated case formulation", 1),
    ("Day 1 recap", None, 1),
    ("Research-informed assessment practice", None, 1),
    ("Reading List", None, 1),
    ("Exercise Instructions (Retained from Earlier Drafts)", None, 1),
]


def slide_text(slide) -> str:
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    )


def find_slide(prs: Presentation, needle: str, start: int = 0) -> int | None:
    needle_l = needle.lower()
    for i in range(start, len(prs.slides)):
        if needle_l in slide_text(prs.slides[i]).lower():
            return i
    return None


def style_content_title_body(slide) -> None:
    t = slide.shapes.title
    if t.has_text_frame:
        tf = t.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            p.font.name = TITLE_FONT
            p.font.size = Pt(TITLE_PT)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = TITLE_FONT
                r.font.size = Pt(TITLE_PT)
                r.font.bold = True
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 or not ph.has_text_frame:
            continue
        tf = ph.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            p.font.name = BODY_FONT
            p.font.size = Pt(BODY_PT)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = BODY_FONT
                r.font.size = Pt(BODY_PT)
                r.font.bold = True


def set_body_bullets(slide, bullets: list[str]) -> None:
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    if body is None:
        return
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0


def set_title(slide, title: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title


def replace_in_slide_text(slide, old: str, new: str) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or old not in shape.text:
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            if old in p.text:
                p.text = p.text.replace(old, new)
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)


def set_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.strip()
    p.font.name = "Calibri"
    p.font.size = Pt(NOTES_PT)


def append_notes(slide, extra: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    existing = tf.text.strip()
    tf.text = (existing + "\n\n" + extra.strip()) if existing else extra.strip()
    for p in tf.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(NOTES_PT)


def delete_slide(prs: Presentation, index: int) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def add_content_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    set_body_bullets(slide, bullets)
    style_content_title_body(slide)
    set_notes(slide, notes)


def set_exercise_slide(
    prs: Presentation,
    body_needle: str,
    display_lines: list[str],
    notes_extra: str,
    timing_note: str | None = None,
) -> None:
    idx = find_slide(prs, body_needle)
    if idx is None:
        return
    slide = prs.slides[idx]
    set_title(slide, "Session Exercise")
    set_body_bullets(slide, display_lines)
    style_content_title_body(slide)
    note_text = notes_extra
    if timing_note:
        note_text = timing_note + "\n\n" + note_text
    append_notes(slide, note_text)


def build_content() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    shutil.copy2(SOURCE, WORK)
    prs = Presentation(str(WORK))

    append_notes(prs.slides[0], VERSION_NOTE)
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text") and "June 2025" in shape.text:
            shape.text = shape.text.replace("June 2025", "June 2026")

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and "June 2025 — Sunday" in shape.text:
                shape.text = shape.text.replace("June 2025", "June 2026")

    sched_idx = find_slide(prs, "Workshop schedule")
    if sched_idx is not None:
        replace_in_slide_text(
            prs.slides[sched_idx],
            "Four sessions; one exercise per session (20 / 40 / 45 / 60 min)",
            "Four sessions; one exercise per session (10–15 / 25–30 / 15–20 / 35–45 min)",
        )
        append_notes(prs.slides[sched_idx], "Revised V4.7 pacing: lighter mornings.")

    aims_idx = find_slide(prs, "Learning aims")
    if aims_idx is not None:
        append_notes(
            prs.slides[aims_idx],
            "V4.7 LO alignment: egg diagram; love/will; family systems; models integration; "
            "research-informed practice; mild/moderate vs severe disturbance; spiritual crisis framing.",
        )

    for i in range(len(prs.slides) - 1, -1, -1):
        text = slide_text(prs.slides[i])
        if "Mental Health Labels" in text and "pros and cons" in text.lower():
            icd_idx = find_slide(prs, "ICD-11 / DSM-5")
            if icd_idx is not None:
                append_notes(
                    prs.slides[icd_idx],
                    "Discussion prompt (retained): pros and cons of mental health labels.",
                )
            delete_slide(prs, i)
        elif "hopes and fears" in text.lower():
            mild_idx = find_slide(prs, "Understanding mild to moderate")
            if mild_idx is not None:
                append_notes(
                    prs.slides[mild_idx],
                    "Optional opening (retained): hopes and fears about neurosis, psychosis, "
                    "spiritual emergencies, and mental illness.",
                )
            delete_slide(prs, i)
        elif "Definition:" in text and "Defense mechanisms" in text:
            delete_slide(prs, i)
        elif "Question re: Client suitability" in text:
            ind_idx = find_slide(prs, "Indicators of client suitability")
            if ind_idx is not None:
                append_notes(
                    prs.slides[ind_idx],
                    "Pair discussion (retained): suitability factors and exclusion criteria.",
                )
            delete_slide(prs, i)

    # Reload after deletions so new slide parts get unique names in the package
    prs.save(str(WORK))
    prs = Presentation(str(WORK))

    def_open = find_slide(prs, "Defense mechanisms")
    if def_open is not None:
        append_notes(
            prs.slides[def_open],
            "Brief definition: unconscious strategies to manage anxiety and threat.",
        )

    # New slides appended at end — reordered later via COM
    add_content_slide(
        prs,
        "Psychosynthesis egg diagram as an assessment map",
        [
            "Use the egg diagram to hypothesise where distress is organised",
            "Lower unconscious: trauma, fear, primitive defence, unmet need",
            "Middle unconscious: current roles, adaptation, coping style",
            "Higher unconscious: meaning, values, aspiration, emerging potential",
            "Ask: what is wounded, defended, split off, and trying to emerge?",
        ],
        "Primal wounding, splitting, survival personality; developmental/purposive lens.",
    )
    add_content_slide(
        prs,
        "Love and will in assessment",
        [
            "Disturbance may involve blocked love, damaged trust, or distorted will",
            "Ask how the person relates, chooses, avoids, submits, controls, or collapses",
            "Assess agency without blaming the client",
            "Therapy may support the recovery of relatedness, choice, and inner direction",
        ],
        "Love/will and psychopathology — clinically grounded.",
    )
    add_content_slide(
        prs,
        "Family systems and hidden dynamics",
        [
            "Distress may be shaped by family roles, loyalty binds, secrets, and unspoken rules",
            "The identified client may carry conflict that belongs to a wider system",
            "Ask what is expressed individually and what may be systemic",
            "Psychosynthesis assessment includes both person and relational field",
        ],
        "Group dynamics and systemic disturbance in first assessment.",
    )
    add_content_slide(
        prs,
        "From hypothesis to review",
        [
            "Initial assessment creates a working hypothesis, not a fixed truth",
            "Review the formulation as new material, ruptures, and risks emerge",
            "Adjust pace, strategy, and referral thinking over time",
        ],
        "Bridges assessment to ongoing formulation.",
    )
    add_content_slide(
        prs,
        "Psychosynthesis and other models",
        [
            "DSM / ICD can support communication and referral, but do not replace formulation",
            "CBT can clarify thoughts, beliefs, and maintaining cycles",
            "PTM reframes distress through context, survival, and meaning",
            "Psychosynthesis adds subpersonalities, will, meaning, and emergent potential",
        ],
        "How psychosynthesis differs from and integrates other models.",
    )
    add_content_slide(
        prs,
        "Research-informed assessment practice",
        [
            "Use evidence, clinical judgement, and client context together",
            "Screening tools support practice but do not diagnose by themselves",
            "Different methods answer different questions: quantitative, qualitative, case-based",
            "Stay critical, ethical, and up to date",
        ],
        "Brief intro to research-informed assessment for counselling trainees.",
    )
    add_content_slide(
        prs,
        "Exercise Instructions (Retained from Earlier Drafts)",
        [
            "Full original instructions for all session exercises",
            "See speaker notes for complete facilitator detail",
        ],
        "\n\n".join(EXercises.values()),
    )

    set_exercise_slide(
        prs,
        "Purpose and scope of initial assessment",
        [
            "Title: Purpose and scope of initial assessment",
            "Aim: Clarify purpose and scope of initial assessment",
            "Time: 10–15 minutes",
            "Pairs: discuss Q1–Q2, brief share-back",
        ],
        "Full original instructions:\n" + EXercises["purpose_scope"],
        "Day 1 AM — brief exercise only.",
    )
    set_exercise_slide(
        prs,
        "Rahul",
        [
            "Title: Client suitability — case 'Rahul'",
            "Aim: Apply suitability and psychosynthesis lens to a complex referral",
            "Time: 25–30 minutes",
            "Three groups: assess suitability, risks, referral; 2-minute rationale",
        ],
        "Full original instructions:\n" + EXercises["rahul"],
        "Day 1 PM — medium exercise.",
    )
    set_exercise_slide(
        prs,
        "Transference, countertransference, and risk",
        [
            "Title: Transference, countertransference, and risk",
            "Aim: Read relational dynamics and risk in simulated assessment",
            "Time: 15–20 minutes",
            "Groups: scenarios — transference, countertransference, risk; brief debrief",
        ],
        "Full original instructions:\n" + EXercises["transference_risk"],
        "Day 2 AM — brief-to-medium exercise.",
    )
    set_exercise_slide(
        prs,
        "Integrated case formulation workshop",
        [
            "Title: Integrated case formulation workshop",
            "Aim: Psychosynthesis-informed formulation and assessment plan",
            "Time: 35–45 minutes",
            "Groups: Ellen or Kim; formulation map, risk, screening; 5-min presentation",
        ],
        "Full original instructions:\n" + EXercises["formulation"],
        "Day 2 PM — substantial integrative exercise.",
    )

    def_idx = find_slide(prs, "go-to defense mechanisms")
    if def_idx is not None:
        slide = prs.slides[def_idx]
        set_title(slide, "Brief reflection (optional)")
        set_body_bullets(
            slide,
            [
                "What defences do you notice in yourself under stress?",
                "Optional pair prompt — 3–5 minutes",
                "Not a formal session exercise",
            ],
        )
        style_content_title_body(slide)
        append_notes(slide, "Full original instructions:\n" + EXercises["defence_partner"])

    for needle, note in [
        ("Day 1 — Saturday Afternoon", "Exercise 2 — 25–30 minutes."),
        ("Day 2 — Sunday Morning", "Exercise 3 — 15–20 minutes."),
        ("Day 2 — Sunday Afternoon", "Exercise 4 — 35–45 minutes."),
    ]:
        idx = find_slide(prs, needle)
        if idx is not None:
            set_notes(prs.slides[idx], f"- {note}")

    recap_idx = find_slide(prs, "Day 1 recap")
    if recap_idx is not None:
        append_notes(prs.slides[recap_idx], VERSION_NOTE)
    reading_idx = find_slide(prs, "Reading List")
    if reading_idx is not None:
        append_notes(prs.slides[reading_idx], VERSION_NOTE)

    prs.save(str(WORK))
    print(f"Content pass: {len(prs.slides)} slides -> {WORK}")


def com_slide_text(slide) -> str:
    parts = []
    for i in range(1, slide.Shapes.Count + 1):
        sh = slide.Shapes(i)
        if sh.HasTextFrame:
            parts.append(sh.TextFrame.TextRange.Text.strip())
    return "\n".join(p for p in parts if p)


def com_match(slide, title_part: str, body_part: str | None, occurrence: int) -> bool:
    text = com_slide_text(slide)
    if title_part.lower() not in text.lower():
        return False
    if body_part and body_part.lower() not in text.lower():
        return False
    return True


def reorder_with_com() -> None:
    import win32com.client

    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True
    pres = ppt.Presentations.Open(str(WORK.resolve()))
    try:
        count = pres.Slides.Count
        slide_ids = [pres.Slides(i).SlideID for i in range(1, count + 1)]

        matched: list[int] = []
        used: set[int] = set()

        for title_part, body_part, occurrence in FINAL_ORDER:
            hits = 0
            found_id = None
            for idx in range(1, count + 1):
                sid = slide_ids[idx - 1]
                if sid in used:
                    continue
                slide = pres.Slides(idx)
                if com_match(slide, title_part, body_part, occurrence):
                    hits += 1
                    if hits == occurrence:
                        found_id = sid
                        break
            if found_id is None:
                raise RuntimeError(f"Could not match slide: {title_part!r} / {body_part!r} #{occurrence}")
            matched.append(found_id)
            used.add(found_id)

        if len(matched) != count:
            raise RuntimeError(f"Matched {len(matched)} slides but deck has {count}")

        for target_pos, slide_id in enumerate(matched, start=1):
            for idx in range(1, pres.Slides.Count + 1):
                if pres.Slides(idx).SlideID == slide_id:
                    pres.Slides(idx).MoveTo(target_pos)
                    break

        pres.SaveAs(str(OUT.resolve()))
        print(f"Reordered and saved: {OUT} ({pres.Slides.Count} slides)")
    finally:
        pres.Close()
        ppt.Quit()


def verify() -> None:
    prs = Presentation(str(OUT))
    print(f"Verify: {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, 1):
        t = slide.shapes.title.text if slide.shapes.title else "?"
        print(f"  {i:2d}. {t[:70]}")


def main() -> None:
    build_content()
    reorder_with_com()
    verify()
    if WORK.exists():
        WORK.unlink()


if __name__ == "__main__":
    main()
