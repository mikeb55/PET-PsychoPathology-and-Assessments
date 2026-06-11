# -*- coding: utf-8 -*-
"""Build V4.8 Psychological Disorders and their Assessment from V4.7 base."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

PPT_DIR = Path(__file__).resolve().parent
SOURCE = PPT_DIR / "V4.7 Psychological Disorders and their Assessment.pptx"
OUT = PPT_DIR / "V4.8-Psychological-Disorders-and-their-Assessment.pptx"
WORK = PPT_DIR / "_v48_work.pptx"

TITLE_PT = 32
BODY_PT = 22
BODY_FONT = "Arial Narrow"
TITLE_FONT = "Calibri"
NOTES_PT = 11

VERSION_NOTE = (
    "V4.8: one primary case-based exercise per session (Chen / Rahul / scenarios / Ellen–Kim); "
    "progressive difficulty; durations in speaker notes only."
)

EXercises = {
    "chen": """Case study: Chen — PTM application (Day 1 Sat AM)
Duration: 20–25 minutes
Aim: Practise PTM questions and intersectionality thinking with an international student struggling with adjustment/academic performance.
Handout: Chen-PTM-answers.docx (facilitator background — do not paste full text on slides).
Groups apply PTM and intersectionality; discuss basic suitability and what else to explore; brief share-back.""",
    "rahul": """Title: Client suitability — case 'Rahul' (Day 1 Sat PM)
Duration: 35–40 minutes
Aim: Apply suitability criteria, risk assessment, and psychosynthesis lens to a complex referral (bipolar disorder, BPD traits, high-risk behaviours).
Handout: Suitability-for-Counselling-Case-Exercise-Rahul.docx
Three groups with case handout; assess suitability, risks, referral; modality (online vs in-person); 2-minute rationale: work now, wait, or refer.""",
    "transference_risk": """Title: Transference, countertransference, and risk (Day 2 Sun AM)
Duration: 30–40 minutes
Aim: Practise reading transference, countertransference, and risk in 3–5 short assessment scenarios.
Five scenarios in groups; pairs identify transference, countertransference, risk; 5-minute debrief per scenario.
Only case-based/clinical exercise in Sunday AM block.""",
    "formulation": """Title: Integrated case formulation workshop (Day 2 Sun PM — capstone)
Duration: 60 minutes (must not exceed 70)
Aim: Produce a psychosynthesis-informed formulation and assessment plan (Ellen BPD or Kim ED).
Handout: Ellen-BPD-initial-assessment.docx for Ellen (Kim = parallel ED option).
Use formulation map; integrate psychopathology, risk, screening (CORE-OM/CORE-10, PHQ-9, GAD-7, dissociation measures as appropriate); 5-minute group presentation.""",
}

FINAL_ORDER: list[tuple[str, str | None, int]] = [
    ("Psychological Disorders and their Assessment", "Weekend Workshop", 1),
    ("Day 1", "Saturday Morning", 1),
    ("Learning aims", None, 1),
    ("Workshop schedule", None, 1),
    ("Assessing the individual", None, 1),
    ("Many forms of troubling behaviour", None, 1),
    ("Power", "Threat", 1),
    ("Psychological Mindedness", None, 1),
    ("Disidentification as an assessment lens", None, 1),
    ('Observing self / "I"', None, 1),
    ("Psychosynthesis egg diagram as an assessment map", None, 1),
    ("Psychosynthesis assessment questions", None, 1),
    ("Subpersonalities in assessment", None, 1),
    ("Psychosynthesis capacities for therapy", None, 1),
    ("Love and will in assessment", None, 1),
    ("Intersectionality in assessment", None, 1),
    ("Family systems and hidden dynamics", None, 1),
    ("Case study: Chen", None, 1),
    ("Day 1", "Saturday Afternoon", 1),
    ("Online vs in-person assessment", None, 1),
    ("Assessing for suicide and self-harm", None, 1),
    ("Initial screening tools", None, 1),
    ("CORE-OM", None, 1),
    ("CORE-10", None, 1),
    ("PHQ-9", None, 1),
    ("GAD-7", None, 1),
    ("Indicators of client suitability", None, 1),
    ("Higher concern / referral indicators", None, 1),
    ("Defense mechanisms", "QUESTIONS", 1),
    ("Defence mechanisms 1", None, 1),
    ("Defence mechanisms 6", None, 1),
    ("Mental status exam", None, 1),
    ("MSE in practice", None, 1),
    ("Psychosynthesis assessment domains", None, 1),
    ("Session Exercise", "Rahul", 1),
    ("Psychological Disorders and their Assessment", "Sunday", 1),
    ("Day 2", "Sunday Morning", 1),
    ("Opening reflection", None, 1),
    ("Brief reflection (optional)", None, 1),
    ("transference", None, 1),
    ("counter-transference", None, 1),
    ("Transferential field", None, 1),
    ("Dissociation in the transferential field", None, 1),
    ("Risk Assessment", None, 1),
    ("What questions do I ask", None, 1),
    ("Self-harm vs suicidality", None, 1),
    ("Dissociation — screening", None, 1),
    ("Session Exercise", "Transference, countertransference", 1),
    ("Day 2", "Sunday Afternoon", 1),
    ("Understanding mild to moderate", None, 1),
    ("Neurosis, psychosis", None, 1),
    ("Lower, middle, and higher unconscious", None, 1),
    ("Personality disorders", None, 1),
    ("Borderline patterns", None, 1),
    ("Narcissistic patterns", None, 1),
    ("Writing session notes", None, 1),
    ("ICD-11 / DSM-5", None, 1),
    ("Eating disorders", None, 1),
    ("From hypothesis to review", None, 1),
    ("Hypothesising and formulation", None, 1),
    ("Psychosynthesis formulation map", None, 1),
    ("What makes an assessment psychosynthesis-informed", None, 1),
    ("Cognitive tools and negative automatic thoughts", None, 1),
    ("Psychosynthesis and other models", None, 1),
    ("Session Exercise", "Integrated case formulation", 1),
    ("Day 1 recap", None, 1),
    ("Research-informed assessment practice", None, 1),
    ("Reading List", None, 1),
    ("Exercise Instructions", None, 1),
]


def slide_text(slide) -> str:
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    )


def find_slide(prs: Presentation, needle: str, start: int = 0) -> int | None:
    needle_l = needle.lower()
    for i in range(start, len(prs.slides)):
        if needle_l in slide_text(prs.slides[i]).lower():
            return i
    return None


def style_content_title_body(slide) -> None:
    t = slide.shapes.title
    if t.has_text_frame:
        tf = t.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            p.font.name = TITLE_FONT
            p.font.size = Pt(TITLE_PT)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = TITLE_FONT
                r.font.size = Pt(TITLE_PT)
                r.font.bold = True
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 or not ph.has_text_frame:
            continue
        tf = ph.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            p.font.name = BODY_FONT
            p.font.size = Pt(BODY_PT)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = BODY_FONT
                r.font.size = Pt(BODY_PT)
                r.font.bold = True


def set_body_bullets(slide, bullets: list[str]) -> None:
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    if body is None:
        return
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0


def set_title(slide, title: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = title


def replace_in_slide_text(slide, old: str, new: str) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or old not in shape.text:
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            if old in p.text:
                p.text = p.text.replace(old, new)
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)


def set_notes(slide, text: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.strip()
    p.font.name = "Calibri"
    p.font.size = Pt(NOTES_PT)


def append_notes(slide, extra: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    existing = tf.text.strip()
    tf.text = (existing + "\n\n" + extra.strip()) if existing else extra.strip()
    for p in tf.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(NOTES_PT)


def delete_slide(prs: Presentation, index: int) -> None:
    r_id = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def set_exercise_slide(
    prs: Presentation,
    body_needle: str,
    display_lines: list[str],
    notes: str,
) -> None:
    idx = find_slide(prs, body_needle)
    if idx is None:
        raise RuntimeError(f"Exercise slide not found: {body_needle}")
    slide = prs.slides[idx]
    set_title(slide, "Session Exercise")
    set_body_bullets(slide, display_lines)
    style_content_title_body(slide)
    set_notes(slide, notes)


def build_content() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)
    shutil.copy2(SOURCE, WORK)
    prs = Presentation(str(WORK))

    append_notes(prs.slides[0], VERSION_NOTE)

    aims_idx = find_slide(prs, "Learning aims")
    if aims_idx is not None:
        append_notes(
            prs.slides[aims_idx],
            "V4.8: four case-based exercises (Chen → Rahul → scenarios → Ellen/Kim); difficulty increases each session.",
        )

    # Remove non-case Purpose/scope exercise — Chen is the sole AM exercise
    for i in range(len(prs.slides) - 1, -1, -1):
        text = slide_text(prs.slides[i])
        if "Session Exercise" in text and "Purpose and scope" in text:
            delete_slide(prs, i)
            break

    prs.save(str(WORK))
    prs = Presentation(str(WORK))

    # Chen — primary Day 1 AM exercise (anchor slide retained)
    chen_idx = find_slide(prs, "Case study: Chen")
    if chen_idx is not None:
        slide = prs.slides[chen_idx]
        set_body_bullets(
            slide,
            [
                "International student; academic and family pressure",
                "Apply PTM questions in groups (Chen-PTM-answers.docx)",
                "Intersectionality: culture, class, migration, age",
                "Is Chen suitable now? What else to explore?",
            ],
        )
        style_content_title_body(slide)
        set_notes(slide, EXercises["chen"])

    set_exercise_slide(
        prs,
        "case 'Rahul'",
        [
            "Title: Client suitability — case 'Rahul'",
            "Aim: Suitability, risk, modality, psychosynthesis lens",
            "Three groups with case handout; 2-minute rationale",
        ],
        EXercises["rahul"],
    )

    set_exercise_slide(
        prs,
        "Transference, countertransference, and risk",
        [
            "Title: Transference, countertransference, and risk",
            "Aim: Read relational dynamics and risk in short scenarios",
            "Groups: 3–5 scenarios; identify transferential and risk cues; debrief",
        ],
        EXercises["transference_risk"],
    )

    set_exercise_slide(
        prs,
        "Integrated case formulation workshop",
        [
            "Title: Integrated case formulation workshop",
            "Aim: Psychosynthesis-informed formulation and assessment plan",
            "Groups: Ellen (BPD) or Kim (ED); formulation map + screening tools",
            "5-minute group presentation",
        ],
        EXercises["formulation"],
    )

    sched_idx = find_slide(prs, "Workshop schedule")
    if sched_idx is not None:
        slide = prs.slides[sched_idx]
        replace_in_slide_text(
            slide,
            "Four sessions; one exercise per session (10–15 / 25–30 / 15–20 / 35–45 min)",
            "Four sessions; one case exercise per session (Chen 20–25 / Rahul 35–40 / scenarios 30–40 / formulation 60 min)",
        )
        replace_in_slide_text(
            slide,
            "Four sessions; one exercise per session (20 / 40 / 45 / 60 min)",
            "Four sessions; one case exercise per session (Chen 20–25 / Rahul 35–40 / scenarios 30–40 / formulation 60 min)",
        )
        append_notes(
            slide,
            "V4.8: progressive case exercises; full durations in speaker notes on each exercise slide.",
        )

    for needle, note in [
        ("Day 1 — Saturday Morning", "Primary exercise: Chen case (20–25 min) — see Chen slide notes."),
        ("Day 1 — Saturday Afternoon", "Primary exercise: Rahul case (35–40 min) — see Session Exercise notes."),
        ("Day 2 — Sunday Morning", "Primary exercise: transference/countertransference/risk scenarios (30–40 min)."),
        ("Day 2 — Sunday Afternoon", "Primary exercise: integrated formulation workshop (60 min capstone)."),
    ]:
        idx = find_slide(prs, needle)
        if idx is not None:
            append_notes(prs.slides[idx], note)

    appendix_idx = find_slide(prs, "Exercise Instructions")
    if appendix_idx is not None:
        set_notes(prs.slides[appendix_idx], "\n\n".join(EXercises.values()))

    recap_idx = find_slide(prs, "Day 1 recap")
    if recap_idx is not None:
        append_notes(prs.slides[recap_idx], VERSION_NOTE)

    prs.save(str(WORK))
    print(f"Content pass: {len(prs.slides)} slides -> {WORK}")


def com_slide_text(slide) -> str:
    parts = []
    for i in range(1, slide.Shapes.Count + 1):
        sh = slide.Shapes(i)
        if sh.HasTextFrame:
            parts.append(sh.TextFrame.TextRange.Text.strip())
    return "\n".join(p for p in parts if p)


def com_match(slide, title_part: str, body_part: str | None, occurrence: int) -> bool:
    text = com_slide_text(slide)
    if title_part.lower() not in text.lower():
        return False
    if body_part and body_part.lower() not in text.lower():
        return False
    return True


def reorder_with_com() -> None:
    import win32com.client

    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True
    pres = ppt.Presentations.Open(str(WORK.resolve()))
    try:
        count = pres.Slides.Count
        slide_ids = [pres.Slides(i).SlideID for i in range(1, count + 1)]
        matched: list[int] = []
        used: set[int] = set()

        for title_part, body_part, occurrence in FINAL_ORDER:
            hits = 0
            found_id = None
            for idx in range(1, count + 1):
                sid = slide_ids[idx - 1]
                if sid in used:
                    continue
                slide = pres.Slides(idx)
                if com_match(slide, title_part, body_part, occurrence):
                    hits += 1
                    if hits == occurrence:
                        found_id = sid
                        break
            if found_id is None:
                raise RuntimeError(f"Could not match slide: {title_part!r} / {body_part!r} #{occurrence}")
            matched.append(found_id)
            used.add(found_id)

        if len(matched) != count:
            raise RuntimeError(f"Matched {len(matched)} slides but deck has {count}")

        for target_pos, slide_id in enumerate(matched, start=1):
            for idx in range(1, pres.Slides.Count + 1):
                if pres.Slides(idx).SlideID == slide_id:
                    pres.Slides(idx).MoveTo(target_pos)
                    break

        pres.SaveAs(str(OUT.resolve()))
        print(f"Saved: {OUT} ({pres.Slides.Count} slides)")
    finally:
        pres.Close()
        ppt.Quit()


def verify() -> None:
    prs = Presentation(str(OUT))
    print(f"Verify: {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides, 1):
        t = slide.shapes.title.text if slide.shapes.title else "?"
        print(f"  {i:2d}. {t[:72]}")


def main() -> None:
    build_content()
    reorder_with_com()
    verify()
    if WORK.exists():
        WORK.unlink()


if __name__ == "__main__":
    main()
